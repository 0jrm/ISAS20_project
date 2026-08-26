#!/usr/bin/env python3
"""16:9 architecture slides. Diagrams and equations. ≤40 words of prose per slide."""
from __future__ import annotations

from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.backends.backend_pdf import PdfPages
from matplotlib.patches import (
    Circle,
    FancyArrowPatch,
    FancyBboxPatch,
    FancyBboxPatch as FB,
    Polygon,
    Rectangle,
)
import numpy as np
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

HERE = Path(__file__).resolve().parent
FIGS = HERE / "figs"
FIGS.mkdir(exist_ok=True)
OUT_PDF = HERE / "nespreso_v2_da_architecture.pdf"
OUT_PPTX = HERE / "nespreso_v2_da_architecture.pptx"

W, H = 13.333, 7.5  # 16:9
CHROME = True  # titles burned into PNG/PDF; False for editable PPTX overlays

C = dict(
    ink="#1a1a1a",
    mute="#5a5854",
    paper="#ffffff",
    card="#f7f5f0",
    line="#2a2a2a",
    enc="#2d6a4f",
    sat="#c05621",
    fuse="#0f766e",
    mlp="#9d174d",
    mu="#1e3a5f",
    sig="#6b21a8",
    warp="#b45309",
    res="#be123c",
    ice="#e4f0ea",
    sand="#f3ead6",
    blush="#f6e4ea",
    mist="#e4ebf3",
    lilac="#eee4f5",
)


def _rc():
    plt.rcParams.update(
        {
            "font.size": 11,
            "font.family": "DejaVu Sans",
            "text.color": C["ink"],
            "figure.facecolor": C["paper"],
            "savefig.facecolor": C["paper"],
            "pdf.fonttype": 42,
            "mathtext.default": "regular",
        }
    )


def _new():
    fig = plt.figure(figsize=(W, H), facecolor=C["paper"])
    ax = fig.add_axes([0, 0, 1, 1])
    ax.set_xlim(0, W)
    ax.set_ylim(0, H)
    ax.axis("off")
    ax.add_patch(Rectangle((0, 0), W, H, fc=C["paper"], ec="none", zorder=0))
    return fig, ax


def _rbox(ax, x, y, w, h, fc, ec, lw=1.4, r=0.08, z=2):
    ax.add_patch(
        FancyBboxPatch(
            (x, y),
            w,
            h,
            boxstyle=f"round,pad=0.01,rounding_size={r}",
            facecolor=fc,
            edgecolor=ec,
            lw=lw,
            zorder=z,
        )
    )


def _t(ax, x, y, s, fs=11, c=C["ink"], w="normal", ha="left", va="center", fam=None, rot=0, z=4):
    kw = dict(ha=ha, va=va, fontsize=fs, color=c, fontweight=w, rotation=rot, zorder=z)
    if fam:
        kw["fontfamily"] = fam
    ax.text(x, y, s, **kw)


def _arr(ax, a, b, c=C["ink"], lw=1.6, rad=0, ms=12):
    ax.add_patch(
        FancyArrowPatch(
            a,
            b,
            arrowstyle="-|>",
            mutation_scale=ms,
            lw=lw,
            color=c,
            connectionstyle=f"arc3,rad={rad}",
            zorder=5,
        )
    )


def _code(ax, x, y, s, c=C["ink"], fs=9.5):
    _t(ax, x, y, s, fs=fs, c=c, fam="DejaVu Sans Mono")


def _rail(ax, x, y0, y1, color):
    ax.add_patch(Rectangle((x, y0), 0.10, y1 - y0, fc=color, ec="none", zorder=3))


def _cube(ax, x, y, w, h, d, fc, ec=C["line"], lw=1.1):
    """Isometric block (front, top, side)."""
    dx, dy = d * 0.55, d * 0.38
    front = [(x, y), (x + w, y), (x + w, y + h), (x, y + h)]
    top = [(x, y + h), (x + dx, y + h + dy), (x + w + dx, y + h + dy), (x + w, y + h)]
    side = [(x + w, y), (x + w + dx, y + dy), (x + w + dx, y + h + dy), (x + w, y + h)]
    ax.add_patch(Polygon(top, fc=_tint(fc, 1.12), ec=ec, lw=lw, zorder=3))
    ax.add_patch(Polygon(side, fc=_tint(fc, 0.78), ec=ec, lw=lw, zorder=3))
    ax.add_patch(Polygon(front, fc=fc, ec=ec, lw=lw, zorder=4))


def _tint(hexcol, f):
    h = hexcol.lstrip("#")
    r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
    r = int(max(0, min(255, r * f)))
    g = int(max(0, min(255, g * f)))
    b = int(max(0, min(255, b * f)))
    return f"#{r:02x}{g:02x}{b:02x}"


def _footer(ax, n, ntot):
    if not CHROME:
        return
    _t(ax, W - 0.35, 0.22, f"{n} / {ntot}", fs=9, c=C["mute"], ha="right")


def _title(ax, s):
    if not CHROME:
        return
    _t(ax, 0.4, H - 0.38, s, fs=22, w="bold")


def _cap(ax, s):
    if not CHROME:
        return
    _t(ax, 0.4, H - 0.78, s, fs=13, c=C["mute"])


def _words(*parts: str) -> int:
    return sum(len(p.replace("—", " ").split()) for p in parts if p)


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------


def s01_roster(ax):
    title, cap = "v2 models for data assimilation", "Shared encoder. Distinct decode."
    _title(ax, title)
    _cap(ax, cap)
    specs = [
        (0.5, C["ice"], C["enc"], "A_CRPS", "PatchConvMLP", r"$9 \to 32$", "native-z PCA"),
        (4.75, C["sand"], C["warp"], "HeaveFast", "HeaveResidualFast", r"$11 \to 35$", "warp + residual PCA"),
        (9.0, C["lilac"], C["sig"], "ops", "HeaveResidualFast", r"$30 \to 35$", "warp + 19 operators"),
    ]
    for x, fc, ec, name, cls, dim, dec in specs:
        _rbox(ax, x, 1.5, 3.9, 4.6, fc, ec, lw=2.0, r=0.12)
        _t(ax, x + 1.95, 5.45, name, fs=22, w="bold", ha="center", c=ec)
        _t(ax, x + 1.95, 4.7, cls, fs=12, ha="center")
        _t(ax, x + 1.95, 3.7, dim, fs=20, ha="center")
        _t(ax, x + 1.95, 2.5, dec, fs=13, ha="center", c=C["mute"])
    _footer(ax, 1, 12)
    return _words(title, cap), title, cap


def s02_backbone(ax):
    title, cap = "PatchConvMLP", r"patch_shape is None. Satellite encoding is Linear. Conv2d is unused."
    _title(ax, title)
    _cap(ax, cap)

    # nested code
    _rbox(ax, 0.35, 0.55, 5.15, 5.85, "#fbfaf7", C["line"], lw=1.2)
    _rail(ax, 0.5, 5.15, 6.05, C["enc"])
    _rail(ax, 0.5, 4.45, 5.05, C["sat"])
    _rail(ax, 0.5, 3.55, 4.25, C["fuse"])
    _rail(ax, 0.5, 2.15, 3.35, C["mlp"])
    _rail(ax, 0.5, 1.35, 1.95, C["mu"])
    _rail(ax, 0.5, 0.75, 1.25, C["sig"])
    _code(ax, 0.75, 5.75, "enc = x[:, :n_enc]", C["enc"], 10)
    _code(ax, 0.75, 5.35, "h_e = Linear(n_enc, 128)(enc)", C["enc"], 10)
    _code(ax, 0.75, 4.85, "sat = x[:, n_enc:]", C["sat"], 10)
    _code(ax, 0.75, 4.45, "h_s = Linear(n_sat, 128)(sat)", C["sat"], 10)
    _code(ax, 0.75, 3.90, "h = h_e + h_s              # conv is None", C["fuse"], 10)
    _code(ax, 0.75, 3.20, "h = Drop(ReLU(Linear(128,1024)(h)))", C["mlp"], 9)
    _code(ax, 0.75, 2.70, "h = Drop(ReLU(Linear(1024,1024)(h)))", C["mlp"], 9)
    _code(ax, 0.75, 1.65, "μ = Linear(1024, d)(h)", C["mu"], 10)
    _code(ax, 0.75, 1.00, "σ = softplus(Linear(1024,d)(h))+ε", C["sig"], 10)

    # isometric flow
    xs = [6.0, 7.35, 8.7, 10.05, 11.4]
    cols = [C["enc"], C["sat"], C["fuse"], C["mlp"], C["mlp"]]
    labs = ["enc\nLinear", "sat\nLinear", "+", "1024", "1024"]
    for i, (x, fc, lab) in enumerate(zip(xs, cols, labs)):
        _cube(ax, x, 3.35, 0.95, 1.15, 0.55, fc)
        _t(ax, x + 0.48, 3.9, lab, fs=9, ha="center", c="white" if i != 2 else C["ink"], w="bold")
        if i < len(xs) - 1:
            _arr(ax, (x + 1.05, 3.9), (xs[i + 1] - 0.05, 3.9), c=C["mute"], lw=1.3)
    _arr(ax, (6.5, 4.7), (8.95, 4.7), c=C["fuse"], lw=2.0, rad=-0.12)
    _t(ax, 7.7, 5.15, "add", fs=11, c=C["fuse"], w="bold")
    _t(ax, 7.82, 3.05, "Conv2d unused", fs=10, ha="center", c=C["mute"])

    _cube(ax, 10.05, 1.55, 0.95, 0.95, 0.5, C["mu"])
    _cube(ax, 11.4, 1.55, 0.95, 0.95, 0.5, C["sig"])
    _arr(ax, (10.52, 3.3), (10.52, 2.6), c=C["mu"])
    _arr(ax, (11.87, 3.3), (11.87, 2.6), c=C["sig"])
    _t(ax, 10.52, 1.25, r"$\mu$", fs=16, ha="center", c=C["mu"], w="bold")
    _t(ax, 11.87, 1.25, r"$\sigma$", fs=16, ha="center", c=C["sig"], w="bold")

    _t(ax, 8.7, 0.85, r"$h_0=W_e x_{\mathrm{enc}}+W_s x_{\mathrm{sat}}$", fs=13, ha="center")
    _footer(ax, 2, 12)
    return _words(title, cap), title, cap


def s03_inputs(ax):
    title, cap = "Input layout", "Harmonics, optional ENSO, then satellite columns."
    _title(ax, title)
    _cap(ax, cap)

    def row(y, name, chunks, note):
        _t(ax, 0.45, y + 0.55, name, fs=14, w="bold")
        x = 2.6
        for lab, w, fc, ec in chunks:
            _rbox(ax, x, y, w, 0.95, fc, ec, lw=1.5, r=0.06)
            _t(ax, x + w / 2, y + 0.48, lab, fs=12, ha="center", w="bold")
            x += w + 0.08
        _t(ax, 12.95, y + 0.48, note, fs=13, ha="right", c=C["mute"])

    row(
        5.15,
        "A_CRPS",
        [("6 harmonics", 4.2, C["mist"], C["mu"]), ("SSS SST SSH", 3.4, C["sand"], C["sat"])],
        r"$n_{\mathrm{enc}}=6,\ n_{\mathrm{sat}}=3$",
    )
    row(
        3.55,
        "HeaveFast",
        [
            ("6 harmonics", 3.3, C["mist"], C["mu"]),
            ("ONI RONI", 2.1, C["ice"], C["enc"]),
            ("SSS SST SSH", 3.2, C["sand"], C["sat"]),
        ],
        r"$n_{\mathrm{enc}}=8,\ n_{\mathrm{sat}}=3$",
    )
    row(
        1.95,
        "ops",
        [
            ("6 harm.", 1.8, C["mist"], C["mu"]),
            ("ONI RONI", 1.7, C["ice"], C["enc"]),
            ("SSS SST SSH", 2.3, C["sand"], C["sat"]),
            ("19 operators", 4.6, C["lilac"], C["sig"]),
        ],
        r"$n_{\mathrm{enc}}=11,\ n_{\mathrm{sat}}=19$",
    )
    _t(ax, 0.45, 1.15, r"$(\cos,\sin)\,2\pi\,\mathrm{doy}/365,\quad (\cos,\sin)\,2\pi\,\mathrm{lat}/180,\quad (\cos,\sin)\,2\pi\,\mathrm{lon}/360$", fs=13)
    _footer(ax, 3, 12)
    return _words(title, cap), title, cap


def s04_head(ax):
    title, cap = "Heteroscedastic head", "Only μ is assimilated."
    _title(ax, title)
    _cap(ax, cap)
    _rbox(ax, 4.15, 4.55, 5.0, 1.55, C["blush"], C["mlp"], lw=2)
    _t(ax, 6.65, 5.32, r"trunk  $h\in R^{1024}$", fs=18, ha="center", w="bold", c=C["mlp"])
    _arr(ax, (5.4, 4.5), (3.4, 3.35), c=C["mu"], lw=2)
    _arr(ax, (7.9, 4.5), (9.9, 3.35), c=C["sig"], lw=2)
    _rbox(ax, 1.3, 1.55, 4.2, 1.7, C["mist"], C["mu"], lw=2)
    _t(ax, 3.4, 2.7, r"$\mu=W_\mu h+b_\mu$", fs=18, ha="center", c=C["mu"], w="bold")
    _t(ax, 3.4, 2.05, "unconstrained", fs=12, ha="center", c=C["mute"])
    _rbox(ax, 7.85, 1.55, 4.2, 1.7, C["lilac"], C["sig"], lw=2)
    _t(ax, 9.95, 2.7, r"$\sigma=\mathrm{softplus}(W_\sigma h)+\varepsilon$", fs=16, ha="center", c=C["sig"], w="bold")
    _t(ax, 9.95, 2.05, r"$\varepsilon=10^{-3}$", fs=12, ha="center", c=C["mute"])
    _footer(ax, 4, 12)
    return _words(title, cap), title, cap


def s05_acrps(ax):
    title, cap = "A_CRPS reconstruction", "Native 1 m grid. 16 T and 16 S principal components."
    _title(ax, title)
    _cap(ax, cap)

    z = np.linspace(0, 1, 120)
    mean = 18 + 10 * np.exp(-((z - 0.12) / 0.08) ** 2) - 8 * z
    m1 = 2.2 * np.exp(-((z - 0.18) / 0.07) ** 2) * np.sin(2 * np.pi * z)
    rec = mean + 0.8 * m1

    def panel(x0, curve, col, lab, eq=None):
        ax.add_patch(Rectangle((x0, 1.15), 3.5, 4.55, fc="#fafafa", ec=C["line"], lw=1.1, zorder=2))
        xs = x0 + 0.45 + (curve - curve.min()) / (curve.max() - curve.min() + 1e-9) * 2.6
        ys = 5.4 - z * 3.9
        ax.plot(xs, ys, color=col, lw=2.4, zorder=4)
        _t(ax, x0 + 1.75, 0.85, lab, fs=13, ha="center", w="bold", c=col)
        ax.plot([x0 + 0.35, x0 + 0.35], [1.4, 5.45], color="#bbbbbb", lw=0.8, zorder=3)

    panel(0.4, mean, C["mute"], r"$\bar T(z)$")
    _t(ax, 4.15, 3.5, "+", fs=28, ha="center", w="bold")
    panel(4.55, 0.5 + m1, C["mlp"], r"$v_k^{(T)}(z)$")
    _t(ax, 8.3, 3.5, r"$\times\alpha_k$", fs=16, ha="center")
    _t(ax, 9.05, 3.5, "=", fs=28, ha="center", w="bold")
    panel(9.45, rec, C["enc"], r"$T(z)$")

    _t(
        ax,
        6.66,
        6.35,
        r"$T(z)=\bar T(z)+\sum_{k=1}^{16}\alpha_k v_k^{(T)}(z)$",
        fs=16,
        ha="center",
    )
    _footer(ax, 5, 12)
    return _words(title, cap), title, cap


def s06_crps(ax):
    title, cap = "CRPS", "Independent Gaussian scores. No predicted covariance."
    _title(ax, title)
    _cap(ax, cap)
    _rbox(ax, 0.45, 3.55, 7.7, 2.55, C["mist"], C["sig"], lw=1.6)
    _t(
        ax,
        4.3,
        5.35,
        r"$\mathrm{CRPS}=\sigma[z(2\Phi(z)-1)+2\phi(z)-\pi^{-1/2}]$",
        fs=16,
        ha="center",
        c=C["sig"],
    )
    _t(ax, 4.3, 4.35, r"$z=(y-\mu)/\sigma$", fs=16, ha="center")

    xx = np.linspace(-3.4, 3.4, 250)
    pdf = np.exp(-0.5 * xx**2) / np.sqrt(2 * np.pi)
    bx, by, bw, bh = 8.5, 1.3, 4.4, 4.8
    ax.add_patch(Rectangle((bx, by), bw, bh, fc="#fafafa", ec=C["line"], lw=1.1, zorder=2))
    xs = bx + 0.35 + (xx - xx.min()) / (xx.max() - xx.min()) * (bw - 0.7)
    ys = by + 0.5 + pdf / pdf.max() * (bh - 1.3)
    ax.plot(xs, ys, color=C["sig"], lw=2.6, zorder=4)
    xt = bx + 0.35 + (0.85 - xx.min()) / (xx.max() - xx.min()) * (bw - 0.7)
    ax.plot([xt, xt], [by + 0.45, by + bh - 0.55], color=C["res"], lw=2.0, ls="--", zorder=4)
    _t(ax, xt + 0.2, by + bh - 0.4, r"$y$", fs=14, c=C["res"])
    _t(ax, bx + bw / 2, by + 0.28, r"$\mathcal{N}(\mu,\sigma^2)$", fs=13, ha="center", c=C["sig"])
    _footer(ax, 6, 12)
    return _words(title, cap), title, cap


def s07_warp(ax):
    title, cap = "Heave warp", "Piecewise-linear map of MLD and D26 onto a canonical grid."
    _title(ax, title)
    _cap(ax, cap)

    def axis(x, y, h, knots, labels, col, title_s):
        ax.plot([x, x], [y, y + h], color=col, lw=3.2, zorder=3, solid_capstyle="round")
        _t(ax, x, y + h + 0.28, title_s, fs=14, ha="center", w="bold", c=col)
        for k, lab in zip(knots, labels):
            yy = y + h * (1 - k)
            ax.plot([x - 0.12, x + 0.12], [yy, yy], color=col, lw=2.2, zorder=4)
            ax.scatter([x], [yy], s=36, color=col, zorder=5)
            _t(ax, x + 0.28, yy, lab, fs=13, c=col, va="center")

    # Physical knots sit shallower than the template, so the map stretches.
    axis(2.4, 1.05, 5.0, [0, 0.12, 0.30, 1.0], ["0 m", "MLD", "D26", "1800 m"], C["sat"], "physical $z$")
    axis(8.3, 1.05, 5.0, [0, 0.22, 0.42, 1.0], ["0 m", "50 m", "120 m", "1800 m"], C["enc"], "canonical $z$")
    pairs = [(0.0, 0.0), (0.12, 0.22), (0.30, 0.42), (1.0, 1.0)]
    y0, h = 1.05, 5.0
    for a, b in pairs:
        _arr(
            ax,
            (2.55, y0 + h * (1 - a)),
            (8.15, y0 + h * (1 - b)),
            c=C["warp"],
            lw=1.8,
        )
    # trapezoid bands follow the sloped map
    bands = [
        (0.0, 0.12, 0.0, 0.22, "#fde7d0"),
        (0.12, 0.30, 0.22, 0.42, "#e7f3ee"),
        (0.30, 1.0, 0.42, 1.0, "#e7eef6"),
    ]
    for a0, a1, b0, b1, fc in bands:
        ax.add_patch(
            Polygon(
                [
                    (2.55, y0 + h * (1 - a0)),
                    (8.15, y0 + h * (1 - b0)),
                    (8.15, y0 + h * (1 - b1)),
                    (2.55, y0 + h * (1 - a1)),
                ],
                fc=fc,
                ec="none",
                alpha=0.45,
                zorder=1,
            )
        )
    _t(ax, 5.35, 0.55, r"$\mathrm{MLD}=50\,e^{r_0},\quad \mathrm{D26}=\mathrm{MLD}+5+65\,e^{r_1}$", fs=15, ha="center")
    _footer(ax, 7, 12)
    return _words(title, cap), title, cap


def s08_residual(ax):
    title, cap = "Heave residual", "Warp climatology, add residual PCA, invert the map."
    _title(ax, title)
    _cap(ax, cap)
    steps = [
        (0.4, C["sand"], C["warp"], "1", r"$T_{\mathrm{clim}}(z)$", "basin mean"),
        (3.55, C["ice"], C["enc"], "2", r"warp$(T_{\mathrm{clim}})$", "+ residual PCA"),
        (6.7, C["blush"], C["res"], "3", r"$T_c + V\alpha$", "canonical grid"),
        (9.85, C["mist"], C["mu"], "4", r"unwarp$(T_c)$", r"$T(z)$ physical"),
    ]
    z = np.linspace(0, 1, 80)
    shapes = [
        0.35 + 0.45 * np.exp(-((z - 0.38) / 0.10) ** 2),
        0.35 + 0.45 * np.exp(-((z - 0.20) / 0.10) ** 2),
        0.35 + 0.45 * np.exp(-((z - 0.20) / 0.10) ** 2) + 0.08 * np.sin(10 * np.pi * z) * np.exp(-z),
        0.35 + 0.45 * np.exp(-((z - 0.32) / 0.10) ** 2) + 0.06 * np.sin(8 * np.pi * z) * np.exp(-z),
    ]
    for (x, fc, ec, n, eq, sub), curve in zip(steps, shapes):
        _rbox(ax, x, 2.15, 2.9, 3.55, fc, ec, lw=2.0, r=0.1)
        _t(ax, x + 0.22, 5.35, n, fs=18, w="bold", c=ec)
        xs = x + 0.45 + curve * 1.7
        ys = 5.05 - z * 2.05
        ax.plot(xs, ys, color=ec, lw=2.0, zorder=4)
        _t(ax, x + 1.45, 2.55, eq, fs=13, ha="center")
        if x < 9:
            _arr(ax, (x + 3.0, 3.9), (x + 3.45, 3.9), c=C["ink"], lw=2.0)
    _t(ax, 6.66, 1.35, r"$T_{\mathrm{res}}=\alpha_T V_T+\bar m_T\quad$ (train-split warped residuals)", fs=13, ha="center")
    _footer(ax, 8, 12)
    return _words(title, cap), title, cap


def s09_loss(ax):
    title, cap = "Heave loss", "Geometry, CRPS, and a weak 50–200 m Tz term."
    _title(ax, title)
    _cap(ax, cap)
    _t(ax, 6.66, 5.85, r"$L=10\,L_{\mathrm{geom}}+L_{\mathrm{CRPS}}+0.1\,L_{\partial_z T}$", fs=22, ha="center")
    boxes = [
        (0.5, C["sand"], C["warp"], r"$L_{\mathrm{geom}}$", r"MLD$/50$, D26$/120$"),
        (4.7, C["lilac"], C["sig"], r"$L_{\mathrm{CRPS}}$", "MLD, D26, 32 PCs"),
        (8.9, C["blush"], C["res"], r"$L_{\partial_z T}$", "50–200 m"),
    ]
    for x, fc, ec, a, b in boxes:
        _rbox(ax, x, 1.55, 3.85, 3.2, fc, ec, lw=2, r=0.1)
        _t(ax, x + 1.92, 3.7, a, fs=22, ha="center", c=ec, w="bold")
        _t(ax, x + 1.92, 2.55, b, fs=14, ha="center")
    _footer(ax, 9, 12)
    return _words(title, cap), title, cap


def s10_ops(ax):
    title, cap = "ops features", "Nineteen operators sampled from the regional cube."
    _title(ax, title)
    _cap(ax, cap)
    names = [
        r"SST $x,y$",
        r"SST $x,y$ 1°",
        r"SSS $x,y$",
        r"SSS $x,y$ 1°",
        r"SSH $x,y$",
        r"SSH $x,y$ 1°",
        r"$\nabla^2\eta$ 1°",
        r"SST $t$ 7 d",
        r"SSH $t$ 7 d",
        r"$u_g,v_g$",
        r"$u_g,v_g$ 1°",
    ]
    # 19 ops grouped into 11 tiles
    extras = ["local", "1°", "local", "1°", "local", "1°", "lap", "tend", "tend", "local", "1°"]
    # draw 19 small cells in 4 rows
    ops = [
        "sst.gx loc",
        "sst.gy loc",
        "sst.gx 1°",
        "sst.gy 1°",
        "sss.gx loc",
        "sss.gy loc",
        "sss.gx 1°",
        "sss.gy 1°",
        "ssh.gx loc",
        "ssh.gy loc",
        "ssh.gx 1°",
        "ssh.gy 1°",
        r"$\nabla^2\eta$",
        "sst 7d",
        "ssh 7d",
        r"$u_g$ loc",
        r"$v_g$ loc",
        r"$u_g$ 1°",
        r"$v_g$ 1°",
    ]
    cols = 7
    for i, name in enumerate(ops):
        r, c = divmod(i, cols)
        x, y = 0.45 + c * 1.82, 4.85 - r * 1.25
        fc = C["lilac"] if i >= 12 else (C["sand"] if i < 4 or (4 <= i < 8) else C["mist"])
        _rbox(ax, x, y, 1.7, 1.05, fc, C["line"], lw=1.1, r=0.06)
        _t(ax, x + 0.85, y + 0.52, name, fs=11, ha="center")
    _t(ax, 6.66, 0.7, r"$u_g=-(g/f)\partial_y\eta,\quad v_g=(g/f)\partial_x\eta$", fs=15, ha="center")
    _footer(ax, 10, 12)
    return _words(title, cap), title, cap


def s11_serve(ax):
    title, cap = "Assimilation products", "Profiles from μ. Diagonal R from Dai σ_o after H."
    _title(ax, title)
    _cap(ax, cap)
    nodes = [
        (0.5, "inputs", C["mist"], C["mu"]),
        (3.3, "encoder", C["ice"], C["enc"]),
        (6.1, r"$\mu$", C["blush"], C["mlp"]),
        (8.9, "decode", C["sand"], C["warp"]),
        (11.5, r"$T(z),S(z)$", C["ice"], C["enc"]),
    ]
    for x, lab, fc, ec in nodes:
        _rbox(ax, x, 3.55, 2.15, 1.7, fc, ec, lw=2, r=0.1)
        _t(ax, x + 1.07, 4.4, lab, fs=16, ha="center", w="bold", c=ec)
        if x < 11:
            _arr(ax, (x + 2.2, 4.4), (x + 2.55, 4.4), c=C["ink"], lw=2)
    _rbox(ax, 6.1, 1.15, 4.95, 1.55, C["lilac"], C["sig"], lw=2, r=0.1)
    _t(ax, 8.57, 1.92, r"$R=\mathrm{diag}(\sigma_o^{\mathrm{Dai}})\ \mathrm{after}\ H$", fs=16, ha="center", c=C["sig"])
    _arr(ax, (7.17, 3.5), (7.17, 2.8), c=C["sig"], lw=1.6)
    _t(ax, 7.35, 3.15, r"not $\sigma_{\mathrm{CRPS}}$", fs=11, c=C["sig"])
    _footer(ax, 11, 12)
    return _words(title, cap), title, cap


def s12_compare(ax):
    title, cap = "Two reconstructors", "A_CRPS is additive in z. Heave is additive after a vertical map."
    _title(ax, title)
    _cap(ax, cap)
    z = np.linspace(0, 1, 100)
    mean = 0.35 + 0.5 * np.exp(-((z - 0.2) / 0.1) ** 2)
    rec = mean + 0.08 * np.sin(8 * np.pi * z) * np.exp(-2 * z)
    clim = 0.4 + 0.45 * np.exp(-((z - 0.38) / 0.09) ** 2)
    warped = 0.4 + 0.45 * np.exp(-((z - 0.2) / 0.09) ** 2)

    def col(x, fc, ec, head, curves):
        _rbox(ax, x, 0.7, 5.9, 5.55, fc, ec, lw=2.2, r=0.1)
        _t(ax, x + 2.95, 5.85, head, fs=18, ha="center", w="bold", c=ec)
        for i, (curve, colr, lab, xoff) in enumerate(curves):
            xs = x + 0.45 + xoff + curve * 1.35
            ys = 5.35 - z * 3.7
            ax.plot(xs, ys, color=colr, lw=2.3, zorder=4)
            _t(ax, x + 0.9 + xoff + 0.5, 1.15, lab, fs=11, ha="center", c=colr)

    col(
        0.4,
        C["ice"],
        C["enc"],
        "A_CRPS",
        [(mean, C["mute"], r"$\bar T$", 0.2), (rec, C["enc"], r"$T(z)$", 2.4)],
    )
    col(
        7.0,
        C["sand"],
        C["warp"],
        "Heave",
        [(clim, C["mute"], "clim", 0.15), (warped, C["warp"], "warped", 2.35)],
    )
    _footer(ax, 12, 12)
    return _words(title, cap), title, cap


SLIDES = [
    s01_roster,
    s02_backbone,
    s03_inputs,
    s04_head,
    s05_acrps,
    s06_crps,
    s07_warp,
    s08_residual,
    s09_loss,
    s10_ops,
    s11_serve,
    s12_compare,
]

NOTES = [
    "Three served configurations share PatchConvMLP in point mode. They differ in input width and in the map from μ to T(z), S(z).",
    "The class name is historical. DA cells set patch_shape=None, so sat_proj is Linear(n_sat, 128) and self.conv is None. Conv2d is only built for the unused conv3 / patch ablation.",
    "Base encoding is six Fourier features of doy, latitude, and longitude. HeaveFast and ops insert ONI and RONI after those six. ops places local SSS/SST/SSH in the encoder half (n_enc=11).",
    "Forward returns concat(μ, σ). softplus plus 1e-3 enforces σ>0. The analysis uses only μ. Observation error is not this σ.",
    "A_CRPS inverts a 16-component PCA for temperature and a 16-component PCA for salinity on the native 1 m grid, z=0…1800 m. Pair the checkpoint with cache 3adcff404b0b.",
    "Training uses the closed-form Gaussian CRPS, independently per output. Predicted cross-level covariances are not used. Ingest R is diagonal Dai σ_o after H.",
    "Heave predicts MLD and D26 as exponentials of unconstrained logits. Physical knots (0, MLD, D26, 1800 m) map piecewise-linearly onto (0, 50, 120, 1800 m).",
    "Climatology is the basin nanmean of cache profiles. Residual PCA is fit on train-split warped residuals, not the native-z cache PCA. Unwarp returns profiles to physical z.",
    "Geometry loss is dimensionless MSE on MLD/50 and D26/120. CRPS is applied to MLD, D26, and 32 residual scores (34 scalars). Stretch is unused. A weak T_z penalty acts on 50–200 m.",
    "ops appends 19 cube operators: gradients of SST, SSS, SSH at native and 1° scales, SSH Laplacian, 7-day tendencies, and geostrophic velocities. Sample from gom_cube.zarr at request time.",
    "API decode writes T(z), S(z) from μ. Do not write CRPS σ into R. H is applied by TSIS on cycle layer thickness. Do not apply HYCOM layer_sample in the API.",
    "A_CRPS adds PCA modes at fixed depth. Heave first registers two landmarks, adds residual shape on the aligned grid, then inverts the map.",
]


def _pptx_box(slide, l, t, w, h, text, *, size, bold=False, color=(0x1A, 0x1A, 0x1A), align=PP_ALIGN.LEFT):
    shape = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = RGBColor(*color)
    run.font.name = "Calibri"
    return shape


def _render(png_prefix: str):
    pngs = []
    meta = []
    counts = []
    for i, fn in enumerate(SLIDES, 1):
        fig, ax = _new()
        nwords, title, cap = fn(ax)
        counts.append((fn.__name__, nwords))
        if nwords > 40:
            raise SystemExit(f"{fn.__name__} has {nwords} words (max 40)")
        p = FIGS / f"{png_prefix}{i:02d}_{fn.__name__[4:]}.png"
        fig.savefig(p, dpi=160)
        pngs.append(p)
        meta.append((title, cap))
        plt.close(fig)
        print(f"{png_prefix}{i:02d}  {nwords:2d}w  {fn.__name__}")
    return pngs, meta, counts


def main() -> None:
    global CHROME
    _rc()
    CHROME = True
    with PdfPages(OUT_PDF) as pdf:
        d = pdf.infodict()
        d["Title"] = "NeSPReSO v2 architecture"
        d["Author"] = "ISAS20_project"
        for i, fn in enumerate(SLIDES, 1):
            fig, ax = _new()
            nwords, _, _ = fn(ax)
            if nwords > 40:
                raise SystemExit(f"{fn.__name__} has {nwords} words (max 40)")
            pdf.savefig(fig)
            fig.savefig(FIGS / f"{i:02d}_{fn.__name__[4:]}.png", dpi=160)
            plt.close(fig)
            print(f"{i:02d}  {nwords:2d}w  {fn.__name__}")

    CHROME = False
    pngs, meta, _ = _render("pptx_")

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    blank = prs.slide_layouts[6]
    ntot = len(pngs)
    for n, (p, (title, cap), note) in enumerate(zip(pngs, meta, NOTES), 1):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(p), Emu(0), Emu(0), width=Inches(W), height=Inches(H))
        _pptx_box(slide, 0.4, 0.16, 12.2, 0.48, title, size=26, bold=True)
        _pptx_box(slide, 0.4, 0.58, 12.2, 0.38, cap, size=15, color=(0x5A, 0x58, 0x54))
        _pptx_box(
            slide, 11.6, 7.12, 1.4, 0.28, f"{n} / {ntot}", size=11, color=(0x5A, 0x58, 0x54), align=PP_ALIGN.RIGHT
        )
        slide.notes_slide.notes_text_frame.text = note
    prs.save(OUT_PPTX)
    print("wrote", OUT_PDF)
    print("wrote", OUT_PPTX)


if __name__ == "__main__":
    main()
